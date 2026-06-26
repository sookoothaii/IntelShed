"""Briefing & feed-status document export (PDF, DOCX, XLSX).

Uses reportlab (PDF), python-docx (DOCX), openpyxl (XLSX) — all pure-Python,
no system dependencies.  Lazy imports so the backend runs without these
if document export is not needed.
"""

from __future__ import annotations

import io
import json
from datetime import datetime, timezone
from typing import Any

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

_WB_DARK = "#0a0e14"
_WB_ACCENT = "#00ffa3"
_WB_AMBER = "#ffb454"
_WB_RED = "#ff4757"
_WB_GREY = "#5a6470"


def _ts() -> str:
    return datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")


def _safe_text(val: Any, fallback: str = "—") -> str:
    if val is None:
        return fallback
    return str(val)


def _digest_lines(briefing: dict) -> dict[str, list[str]]:
    """Extract digest section lines from a briefing payload."""
    sources = briefing.get("sources") or {}
    digest = sources.get("digest") or {}
    sections = sources.get("_digest_sections") or {}
    return {
        "local": sections.get("local") or [],
        "regional": sections.get("regional") or [],
        "global": sections.get("global") or [],
        "region_label": digest.get("region_label") or "Operator Region",
        "window": digest.get("window") or "24h",
    }


def _briefing_from_db_row(row) -> dict:
    """Parse a briefings DB row into the same shape as GET /api/briefing."""
    sources = {}
    try:
        sources = json.loads(row["sources"]) if row["sources"] else {}
    except Exception:
        pass
    return {
        "created_at": row["created_at"],
        "text": row["text"] or "",
        "sources": sources,
        "quality": sources.get("quality"),
        "watch_items": sources.get("watch_items") or [],
        "insights": sources.get("insights") or [],
        "digest": sources.get("digest") or {},
        "_digest_sections": sources.get("_digest_sections") or {},
    }


# ---------------------------------------------------------------------------
# PDF export (reportlab)
# ---------------------------------------------------------------------------


def briefing_to_pdf(briefing: dict) -> bytes:
    """Render a briefing dict as a formatted PDF report."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.lib.colors import HexColor
    from reportlab.platypus import (
        SimpleDocTemplate,
        Paragraph,
        Spacer,
    )

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        leftMargin=20 * mm,
        rightMargin=20 * mm,
        topMargin=18 * mm,
        bottomMargin=18 * mm,
    )

    styles = getSampleStyleSheet()
    h1 = ParagraphStyle(
        "WBH1",
        parent=styles["Heading1"],
        fontSize=18,
        textColor=HexColor(_WB_DARK),
        spaceAfter=4,
    )
    h2 = ParagraphStyle(
        "WBH2",
        parent=styles["Heading2"],
        fontSize=13,
        textColor=HexColor(_WB_DARK),
        spaceBefore=10,
        spaceAfter=4,
    )
    meta = ParagraphStyle(
        "WBMeta",
        parent=styles["Normal"],
        fontSize=8,
        textColor=HexColor(_WB_GREY),
        spaceAfter=2,
    )
    body = ParagraphStyle(
        "WBBody", parent=styles["Normal"], fontSize=9, leading=13, spaceAfter=2
    )
    bullet = ParagraphStyle(
        "WBBullet", parent=body, leftIndent=12, bulletIndent=2, spaceAfter=1
    )
    story: list = []

    # --- Header ---
    story.append(Paragraph("WORLDBASE — 24h Situation Briefing", h1))
    created = _safe_text(briefing.get("created_at"), "—")
    story.append(Paragraph(f"Generated: {created}", meta))

    quality = briefing.get("quality") or {}
    q_score = quality.get("score")
    if q_score is not None:
        story.append(Paragraph(f"Quality Score: {q_score:.2f}/1.00", meta))

    sections = _digest_lines(briefing)
    story.append(
        Paragraph(
            f"Region: {sections['region_label']}  |  Window: {sections['window']}", meta
        )
    )
    story.append(Spacer(1, 6))

    # --- Digest sections ---
    for label, key in [
        ("LOCAL", "local"),
        ("REGIONAL", "regional"),
        ("GLOBAL", "global"),
    ]:
        lines = sections.get(key) or []
        story.append(Paragraph(label, h2))
        if not lines:
            story.append(Paragraph("No signals.", body))
        else:
            for line in lines:
                story.append(Paragraph(f"• {_safe_text(line)}", bullet))
        story.append(Spacer(1, 4))

    # --- Watch items ---
    watch = briefing.get("watch_items") or []
    if watch:
        story.append(Paragraph("WATCH ITEMS", h2))
        for w in watch[:10]:
            headline = _safe_text(w.get("headline") or w.get("text"), "—")
            horizon = _safe_text(w.get("horizon"), "")
            story.append(Paragraph(f"• {headline}  ({horizon})", bullet))
        story.append(Spacer(1, 4))

    # --- Insights ---
    insights = briefing.get("insights") or []
    if insights:
        story.append(Paragraph("INSIGHT CARDS", h2))
        for ins in insights[:5]:
            headline = _safe_text(ins.get("headline"), "—")
            so_what = _safe_text(ins.get("so_what"), "")
            score = ins.get("score")
            score_str = f"  [score: {score:.2f}]" if score is not None else ""
            story.append(Paragraph(f"• {headline}{score_str}", bullet))
            if so_what:
                story.append(Paragraph(f"  → {so_what}", body))
        story.append(Spacer(1, 4))

    # --- Full briefing text ---
    text = briefing.get("text") or ""
    if text:
        story.append(Paragraph("FULL BRIEFING TEXT", h2))
        for para in text.split("\n"):
            para = para.strip()
            if not para:
                continue
            story.append(Paragraph(para, body))

    # --- Footer ---
    story.append(Spacer(1, 8))
    story.append(Paragraph(f"WorldBase Document Export — {_ts()}", meta))

    doc.build(story)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# DOCX export (python-docx)
# ---------------------------------------------------------------------------


def briefing_to_docx(briefing: dict) -> bytes:
    """Render a briefing dict as a Word document."""
    from docx import Document
    from docx.shared import Pt, RGBColor

    doc = Document()

    # Title
    doc.add_heading("WORLDBASE — 24h Situation Briefing", level=1)
    created = _safe_text(briefing.get("created_at"), "—")
    p = doc.add_paragraph(f"Generated: {created}")
    p.runs[0].font.size = Pt(8)
    p.runs[0].font.color.rgb = RGBColor(0x5A, 0x64, 0x70)

    quality = briefing.get("quality") or {}
    q_score = quality.get("score")
    if q_score is not None:
        p = doc.add_paragraph(f"Quality Score: {q_score:.2f}/1.00")
        p.runs[0].font.size = Pt(8)
        p.runs[0].font.color.rgb = RGBColor(0x5A, 0x64, 0x70)

    sections = _digest_lines(briefing)
    p = doc.add_paragraph(
        f"Region: {sections['region_label']}  |  Window: {sections['window']}"
    )
    p.runs[0].font.size = Pt(8)
    p.runs[0].font.color.rgb = RGBColor(0x5A, 0x64, 0x70)

    # Digest sections
    for label, key in [
        ("LOCAL", "local"),
        ("REGIONAL", "regional"),
        ("GLOBAL", "global"),
    ]:
        lines = sections.get(key) or []
        doc.add_heading(label, level=2)
        if not lines:
            doc.add_paragraph("No signals.")
        else:
            for line in lines:
                doc.add_paragraph(_safe_text(line), style="List Bullet")

    # Watch items
    watch = briefing.get("watch_items") or []
    if watch:
        doc.add_heading("WATCH ITEMS", level=2)
        for w in watch[:10]:
            headline = _safe_text(w.get("headline") or w.get("text"), "—")
            horizon = _safe_text(w.get("horizon"), "")
            doc.add_paragraph(f"{headline}  ({horizon})", style="List Bullet")

    # Insights
    insights = briefing.get("insights") or []
    if insights:
        doc.add_heading("INSIGHT CARDS", level=2)
        for ins in insights[:5]:
            headline = _safe_text(ins.get("headline"), "—")
            so_what = _safe_text(ins.get("so_what"), "")
            score = ins.get("score")
            score_str = f"  [score: {score:.2f}]" if score is not None else ""
            doc.add_paragraph(f"{headline}{score_str}", style="List Bullet")
            if so_what:
                doc.add_paragraph(f"→ {so_what}")

    # Full briefing text
    text = briefing.get("text") or ""
    if text:
        doc.add_heading("FULL BRIEFING TEXT", level=2)
        for para in text.split("\n"):
            para = para.strip()
            if not para:
                continue
            doc.add_paragraph(para)

    # Footer
    doc.add_paragraph()
    footer = doc.add_paragraph(f"WorldBase Document Export — {_ts()}")
    footer.runs[0].font.size = Pt(8)
    footer.runs[0].font.color.rgb = RGBColor(0x5A, 0x64, 0x70)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PPTX export (python-pptx)
# ---------------------------------------------------------------------------


def briefing_to_pptx(briefing: dict) -> bytes:
    """Render a briefing dict as a PowerPoint situational awareness deck."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    dark = RGBColor(0x0A, 0x0E, 0x14)
    accent = RGBColor(0x00, 0xFF, 0xA3)
    amber = RGBColor(0xFF, 0xB4, 0x54)
    grey = RGBColor(0x5A, 0x64, 0x70)
    white = RGBColor(0xFF, 0xFF, 0xFF)

    def _add_bg(slide, color=dark):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_textbox(
        slide,
        left,
        top,
        width,
        height,
        text,
        font_size=14,
        color=white,
        bold=False,
        align=PP_ALIGN.LEFT,
    ):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        return tf

    sections = _digest_lines(briefing)
    quality = briefing.get("quality") or {}
    q_score = quality.get("score")
    created = _safe_text(briefing.get("created_at"), "—")

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg(slide)
    _add_textbox(
        slide,
        Inches(1),
        Inches(2.5),
        Inches(11),
        Inches(1.2),
        "WORLDBASE",
        44,
        accent,
        True,
        PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        Inches(1),
        Inches(3.5),
        Inches(11),
        Inches(0.8),
        "24h Situation Briefing",
        28,
        white,
        False,
        PP_ALIGN.CENTER,
    )
    meta_parts = [f"Generated: {created}", f"Region: {sections['region_label']}"]
    if q_score is not None:
        meta_parts.append(f"Quality: {q_score:.2f}/1.00")
    _add_textbox(
        slide,
        Inches(1),
        Inches(4.5),
        Inches(11),
        Inches(0.5),
        "  |  ".join(meta_parts),
        12,
        grey,
        False,
        PP_ALIGN.CENTER,
    )

    # --- Slide 2: Digest Summary ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_textbox(
        slide,
        Inches(0.5),
        Inches(0.3),
        Inches(12),
        Inches(0.6),
        "DIGEST SUMMARY",
        24,
        accent,
        True,
    )

    y = Inches(1.2)
    for label, key in [
        ("LOCAL", "local"),
        ("REGIONAL", "regional"),
        ("GLOBAL", "global"),
    ]:
        lines = sections.get(key) or []
        _add_textbox(
            slide, Inches(0.5), y, Inches(3), Inches(0.4), label, 16, amber, True
        )
        y2 = y + Inches(0.5)
        if not lines:
            _add_textbox(
                slide, Inches(0.8), y2, Inches(11), Inches(0.4), "No signals.", 11, grey
            )
            y2 += Inches(0.4)
        else:
            for line in lines[:6]:
                _add_textbox(
                    slide,
                    Inches(0.8),
                    y2,
                    Inches(11),
                    Inches(0.35),
                    f"• {_safe_text(line)}",
                    11,
                    white,
                )
                y2 += Inches(0.35)
        y = y2 + Inches(0.2)

    # --- Slide 3: Watch Items ---
    watch = briefing.get("watch_items") or []
    if watch:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_bg(slide)
        _add_textbox(
            slide,
            Inches(0.5),
            Inches(0.3),
            Inches(12),
            Inches(0.6),
            "WATCH ITEMS",
            24,
            accent,
            True,
        )
        y = Inches(1.2)
        for w in watch[:8]:
            headline = _safe_text(w.get("headline") or w.get("text"), "—")
            horizon = _safe_text(w.get("horizon"), "")
            _add_textbox(
                slide,
                Inches(0.8),
                y,
                Inches(11),
                Inches(0.5),
                f"• {headline}  ({horizon})",
                13,
                white,
            )
            y += Inches(0.55)

    # --- Slide 4: Insight Cards ---
    insights = briefing.get("insights") or []
    if insights:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_bg(slide)
        _add_textbox(
            slide,
            Inches(0.5),
            Inches(0.3),
            Inches(12),
            Inches(0.6),
            "INSIGHT CARDS",
            24,
            accent,
            True,
        )
        y = Inches(1.2)
        for ins in insights[:5]:
            headline = _safe_text(ins.get("headline"), "—")
            so_what = _safe_text(ins.get("so_what"), "")
            score = ins.get("score")
            score_str = f"  [{score:.2f}]" if score is not None else ""
            _add_textbox(
                slide,
                Inches(0.8),
                y,
                Inches(11),
                Inches(0.4),
                f"• {headline}{score_str}",
                13,
                white,
                True,
            )
            y += Inches(0.45)
            if so_what:
                _add_textbox(
                    slide,
                    Inches(1.2),
                    y,
                    Inches(10.5),
                    Inches(0.4),
                    f"→ {so_what}",
                    11,
                    grey,
                )
                y += Inches(0.4)
            y += Inches(0.1)

    # --- Slide 5+: Full briefing text (paginated) ---
    text = briefing.get("text") or ""
    if text:
        paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
        # ~12 lines per slide
        per_slide = 12
        for i in range(0, len(paragraphs), per_slide):
            chunk = paragraphs[i : i + per_slide]
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_bg(slide)
            _add_textbox(
                slide,
                Inches(0.5),
                Inches(0.3),
                Inches(12),
                Inches(0.6),
                "FULL BRIEFING TEXT",
                20,
                accent,
                True,
            )
            y = Inches(1.1)
            for para in chunk:
                _add_textbox(
                    slide, Inches(0.5), y, Inches(12), Inches(0.45), para, 10, white
                )
                y += Inches(0.42)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# XLSX feed status export (openpyxl)
# ---------------------------------------------------------------------------


def feeds_to_xlsx(feed_data: list[dict]) -> bytes:
    """Render feed connector status as an Excel spreadsheet."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment

    wb = Workbook()
    ws = wb.active
    ws.title = "Feed Status"

    headers = [
        "Name",
        "Source",
        "Category",
        "Status",
        "Freshness",
        "Count",
        "Last Fetch",
        "Error",
        "URL",
    ]
    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill("solid", start_color="0A0E14")
    for col, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center")

    for row_idx, feed in enumerate(feed_data, 2):
        ws.cell(row=row_idx, column=1, value=feed.get("name") or "—")
        ws.cell(row=row_idx, column=2, value=feed.get("source") or "—")
        ws.cell(row=row_idx, column=3, value=feed.get("category") or "—")
        ws.cell(row=row_idx, column=4, value=feed.get("status") or "—")
        ws.cell(row=row_idx, column=5, value=feed.get("freshness") or "—")
        ws.cell(row=row_idx, column=6, value=feed.get("count") or 0)
        ws.cell(row=row_idx, column=7, value=feed.get("last_fetch") or "—")
        ws.cell(row=row_idx, column=8, value=feed.get("error") or "")
        ws.cell(row=row_idx, column=9, value=feed.get("url") or "")

        # Color-code freshness
        fresh = str(feed.get("freshness") or "").lower()
        if fresh == "fresh":
            ws.cell(row=row_idx, column=5).fill = PatternFill(
                "solid", start_color="00FFA3"
            )
        elif fresh == "stale":
            ws.cell(row=row_idx, column=5).fill = PatternFill(
                "solid", start_color="FFB454"
            )
        elif fresh == "error":
            ws.cell(row=row_idx, column=5).fill = PatternFill(
                "solid", start_color="FF4757"
            )

    # Auto-width
    for col in range(1, len(headers) + 1):
        max_len = len(str(headers[col - 1]))
        for row in range(2, ws.max_row + 1):
            val = ws.cell(row=row, column=col).value
            if val:
                max_len = max(max_len, min(len(str(val)), 50))
        ws.column_dimensions[chr(64 + col)].width = max_len + 3

    # Summary sheet
    ws2 = wb.create_sheet("Summary")
    total = len(feed_data)
    fresh_n = sum(
        1 for f in feed_data if str(f.get("freshness") or "").lower() == "fresh"
    )
    stale_n = sum(
        1 for f in feed_data if str(f.get("freshness") or "").lower() == "stale"
    )
    error_n = sum(
        1 for f in feed_data if str(f.get("freshness") or "").lower() == "error"
    )
    ws2.cell(row=1, column=1, value="WorldBase Feed Status Summary").font = Font(
        bold=True, size=14
    )
    ws2.cell(row=3, column=1, value="Generated")
    ws2.cell(row=3, column=2, value=_ts())
    ws2.cell(row=4, column=1, value="Total Feeds")
    ws2.cell(row=4, column=2, value=total)
    ws2.cell(row=5, column=1, value="Fresh")
    ws2.cell(row=5, column=2, value=fresh_n)
    ws2.cell(row=6, column=1, value="Stale")
    ws2.cell(row=6, column=2, value=stale_n)
    ws2.cell(row=7, column=1, value="Error")
    ws2.cell(row=7, column=2, value=error_n)
    ws2.column_dimensions["A"].width = 18
    ws2.column_dimensions["B"].width = 30

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()
